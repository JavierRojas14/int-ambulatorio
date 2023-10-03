import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def add_dataframes_to_powerpoint(
    dataframes,
    pptx_filename,
    font_size=12,
    font_family="Arial",
    max_cell_characters=50,
    cell_width=2,
    cell_height=1,
):
    # Create a PowerPoint presentation with 16:9 aspect ratio
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)

    for df in dataframes:
        # Create a new slide for each DataFrame
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Use a blank slide layout

        # Define table position on the slide (fits the entire slide)
        left = Inches(0.5)  # Adjust the left margin as needed
        top = Inches(1)
        width = Inches(cell_width)
        height = Inches(cell_height)

        # Add a table to the slide
        table = slide.shapes.add_table(len(df) + 1, len(df.columns), left, top, width, height).table

        # Set column headers and alignment based on data type
        for col_num, column_name in enumerate(df.columns):
            cell = table.cell(0, col_num)
            cell.text = column_name
            cell.text_frame.paragraphs[0].alignment = (
                PP_ALIGN.LEFT if pd.api.types.is_string_dtype(df[column_name]) else PP_ALIGN.RIGHT
            )
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(font_size)
            cell.text_frame.paragraphs[0].font.name = font_family
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black font color

        # Iterate through the DataFrame rows and columns
        for row_idx, row in enumerate(df.itertuples(index=False)):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx + 1, col_idx)

                # Truncate cell content if it exceeds the character limit
                cell_text = str(value)[:max_cell_characters]

                cell.text = cell_text

                # Set alignment and font properties based on data type
                cell.text_frame.paragraphs[0].alignment = (
                    PP_ALIGN.LEFT
                    if pd.api.types.is_string_dtype(df.iloc[:, col_idx])
                    else PP_ALIGN.RIGHT
                )
                cell.text_frame.paragraphs[0].font.size = Pt(font_size)
                cell.text_frame.paragraphs[0].font.name = font_family

        # Autofit table dimensions based on cell content
        table.allow_autofit = True
        table.autofit = True

    # Save the PowerPoint presentation
    prs.save(pptx_filename)
